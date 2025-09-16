"""Generate PowerPoint presentation for AutomatedAIpowered project using the Presentation_Outline.md.

Usage:
    python scripts/generate_presentation.py 

Outputs:
    AutomatedAIpowered_Project.pptx in project root.

Notes:
    - Requires python-pptx installed.
    - Basic parsing of the outline file: slides separated by '## Slide'.
    - Speaker notes populated from **Speaker Notes:** sections.
"""
from __future__ import annotations
import re
from pathlib import Path
from typing import List, Tuple

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
except ImportError as e:
    raise SystemExit("python-pptx not installed. Run: pip install python-pptx") from e

ROOT = Path(__file__).resolve().parents[1]
OUTLINE_FILE = ROOT / "Presentation_Outline.md"
OUTPUT_FILE = ROOT / "AutomatedAIpowered_Project.pptx"

SLIDE_TITLE_LAYOUT_INDEX = 0  # Default title slide
SLIDE_CONTENT_LAYOUT_INDEX = 1  # Title + content

# Optional theming
TITLE_COLOR = "305496"  # Dark blue
ACCENT_COLOR = "4F81BD"
FONT_NAME = "Segoe UI"


def read_outline() -> str:
    if not OUTLINE_FILE.exists():
        raise FileNotFoundError(f"Outline file not found: {OUTLINE_FILE}")
    return OUTLINE_FILE.read_text(encoding="utf-8")


def split_slides(markdown: str) -> List[str]:
    # Split on lines starting with ## Slide
    parts = re.split(r"(?=^## Slide)" , markdown, flags=re.MULTILINE)
    # Remove header part before first slide
    slides = [p.strip() for p in parts if p.strip().startswith("## Slide")]
    return slides


def parse_slide(block: str) -> Tuple[str, List[str], str]:
    """Return (title, bullets, speaker_notes)."""
    # Title line like: ## Slide 2 – Executive Summary
    first_line = block.splitlines()[0]
    m = re.match(r"## Slide\s+\d+\s+[–-]\s+(.+)", first_line)
    if not m:
        title = first_line.replace("##", "").strip()
    else:
        title = m.group(1).strip()

    body_lines: List[str] = []
    notes_lines: List[str] = []
    in_notes = False

    for line in block.splitlines()[1:]:
        if line.strip().startswith("**Speaker Notes:**"):
            in_notes = True
            note_text = line.split("**Speaker Notes:**", 1)[1].strip().strip('*').strip()
            if note_text:
                notes_lines.append(note_text)
            continue
        if in_notes:
            # Stop notes at slide separator
            if line.startswith("---"):
                break
            notes_lines.append(line)
        else:
            if line.startswith("**") and line.endswith("**"):
                continue
            if line.strip():
                body_lines.append(line)

    # Clean bullet content: keep list markers or plain lines
    bullets: List[str] = []
    for l in body_lines:
        stripped = l.strip()
        if stripped.startswith(("- ", "* ")):
            bullets.append(stripped[2:].strip())
        elif stripped.startswith("1.") or re.match(r"^\d+\.\s", stripped):
            bullets.append(stripped)
        elif stripped.startswith("|") and "|" in stripped[1:]:
            bullets.append(stripped)  # tables kept as raw text
        elif stripped.startswith("```"):
            # Code block start/end markers skip
            continue
        else:
            if stripped:
                bullets.append(stripped)

    speaker_notes = "\n".join([l.rstrip() for l in notes_lines if l.strip()])

    return title, bullets[:20], speaker_notes  # limit excessive bullets


def create_presentation(slides: List[Tuple[str, List[str], str]]):
    prs = Presentation()

    # First slide uses title layout
    for idx, (title, bullets, notes) in enumerate(slides):
        if idx == 0:
            layout = prs.slide_layouts[SLIDE_TITLE_LAYOUT_INDEX]
            slide = prs.slides.add_slide(layout)
            slide.shapes.title.text = title
            if bullets:
                subtitle = slide.placeholders[1]
                subtitle.text = bullets[0][:200]
            if notes:
                slide.notes_slide.notes_text_frame.text = notes
            continue

        layout = prs.slide_layouts[SLIDE_CONTENT_LAYOUT_INDEX]
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = title

        body_shape = slide.shapes.placeholders[1]
        tf = body_shape.text_frame
        tf.clear()
        if bullets:
            first = True
            for b in bullets:
                if first:
                    p = tf.paragraphs[0]
                    p.text = b
                    p.font.name = FONT_NAME
                    first = False
                else:
                    p = tf.add_paragraph()
                    p.text = b
                    p.font.name = FONT_NAME
        else:
            p = tf.paragraphs[0]
            p.text = "(No content)"
            p.font.italic = True
        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    prs.save(OUTPUT_FILE)
    return OUTPUT_FILE


def main():
    md = read_outline()
    slide_blocks = split_slides(md)
    slides = [parse_slide(b) for b in slide_blocks]
    out = create_presentation(slides)
    print(f"Presentation generated: {out}")

if __name__ == "__main__":
    main()
